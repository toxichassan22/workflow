"""
PPTX Export Engine — Tenant-aware.
Generates a PowerPoint file from slide HTML using Playwright screenshots.
"""

import base64
import os
import time
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


_ICON_RE = re.compile(r'[\U0001F000-\U0001FAFF\u2600-\u27BF\uFE0F\u200D]')


def _strip_icons(value):
    return _ICON_RE.sub('', str(value or '')).replace('•', '').strip()


def _hex_to_rgb(hex_color):
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 3:
        hex_color = ''.join(c * 2 for c in hex_color)
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def generate_pptx(slides_data, project_name, branding=None, output_dir=None):
    """
    Generate a PPTX from slide HTML. Each slide is rendered via Playwright
    and embedded as an image for pixel-perfect output.

    Args:
        slides_data: list of slide dicts with 'html' and 'title' keys
        project_name: name for the output file
        branding: tenant branding dict
        output_dir: directory to save the PPTX

    Returns:
        str: path to the generated PPTX file
    """
    if output_dir is None:
        output_dir = os.path.join(os.path.dirname(__file__), '..', 'outputs')
    os.makedirs(output_dir, exist_ok=True)

    # Tenant branding
    font_family = 'The Sans Arabic'
    if branding:
        font_family = branding.get('font_family', font_family)

    # Slide dimensions
    slide_ratio = branding.get('slide_ratio', '16:9') if branding else '16:9'
    if slide_ratio == '4:3':
        slide_w_px, slide_h_px = 1280, 960
        slide_w = Inches(10)
        slide_h = Inches(7.5)
    else:
        slide_w_px, slide_h_px = 1280, 720
        slide_w = Inches(13.333)
        slide_h = Inches(7.5)

    # Build font-face CSS
    project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
    font_faces = ''

    # Check for custom uploaded font
    custom_font_b64 = None
    custom_font_format = None
    if branding and branding.get('font_file_path'):
        font_file_path = branding['font_file_path']
        parts = [p for p in font_file_path.split('/') if p]
        tenant_id = parts[1] if len(parts) >= 2 else None
        if tenant_id:
            font_filename = parts[-1]
            local_font_path = os.path.join(project_root, 'uploads', tenant_id, 'fonts', font_filename)
            if os.path.exists(local_font_path):
                try:
                    with open(local_font_path, 'rb') as f:
                        custom_font_b64 = base64.b64encode(f.read()).decode('ascii')
                    ext = os.path.splitext(local_font_path)[1].lower()
                    custom_font_format = 'truetype' if ext == '.ttf' else 'opentype'
                except Exception:
                    pass

    if custom_font_b64 and custom_font_format:
        for w in range(100, 1000, 100):
            font_faces += f"@font-face {{ font-family:'{font_family}'; src:url('data:font/{custom_font_format};base64,{custom_font_b64}') format('{custom_font_format}'); font-weight:{w}; font-style:normal; font-display:swap; }}\n"
    else:
        for font_file, weights in [
            ('assets/fonts/TheSansArabic-Light.otf', [300, 400, 500, 600]),
            ('assets/fonts/BahijTheSansArabic-Bold.ttf', [700, 800, 900]),
        ]:
            fpath = os.path.join(project_root, font_file)
            if os.path.exists(fpath):
                with open(fpath, 'rb') as f:
                    b64 = base64.b64encode(f.read()).decode('ascii')
                fmt = 'opentype' if font_file.endswith('.otf') else 'truetype'
                for w in weights:
                    font_faces += f"@font-face {{ font-family:'{font_family}'; src:url('data:font/{fmt};base64,{b64}') format('{fmt}'); font-weight:{w}; font-style:normal; font-display:swap; }}\n"

    # Embed logo as base64
    logo_data_uri = ''
    if branding and branding.get('logo_path'):
        parts = [p for p in branding['logo_path'].split('/') if p]
        tenant_id = parts[1] if len(parts) >= 2 else None
        if tenant_id:
            tenant_dir = os.path.join(project_root, 'uploads', tenant_id)
            for ext in ['.png', '.jpg', '.jpeg', '.webp']:
                local_path = os.path.join(tenant_dir, f'logo{ext}')
                if os.path.exists(local_path):
                    try:
                        mime = 'image/png' if ext == '.png' else ('image/webp' if ext == '.webp' else 'image/jpeg')
                        with open(local_path, 'rb') as f:
                            b64 = base64.b64encode(f.read()).decode('ascii')
                            logo_data_uri = f'data:{mime};base64,{b64}'
                        break
                    except Exception:
                        pass

    # Render each slide HTML to an image via Playwright
    slide_images = []
    try:
        from playwright.sync_api import sync_playwright
        with sync_playwright() as p:
            browser = p.chromium.launch(args=['--no-sandbox', '--disable-dev-shm-usage', '--disable-gpu'])
            page = browser.new_page(viewport={'width': slide_w_px, 'height': slide_h_px})

            for i, slide_data in enumerate(slides_data):
                html = slide_data.get('html', '')
                if not html:
                    continue
                # Replace logo placeholder
                if logo_data_uri:
                    html = html.replace('##LOGO##', logo_data_uri)
                    if branding and branding.get('logo_path'):
                        html = html.replace(branding['logo_path'], logo_data_uri)
                    html = html.replace('/assets/logo.png', logo_data_uri)

                full_html = f"""<!DOCTYPE html>
<html lang="ar" dir="rtl">
<head>
<meta charset="UTF-8">
<style>
* {{ margin:0; padding:0; box-sizing:border-box; }}
{font_faces}
body {{ direction:rtl; font-family:'{font_family}',Tahoma,Arial,sans-serif; }}
.slide {{ width:{slide_w_px}px; height:{slide_h_px}px; direction:rtl; position:relative; overflow:hidden; }}
img {{ max-width:100%; max-height:100%; object-fit:cover; }}
</style>
</head>
<body>{html}</body>
</html>"""
                page.set_content(full_html, wait_until='load')
                page.wait_for_timeout(500)
                page.evaluate('document.fonts.ready')
                # Screenshot the .slide element
                slide_el = page.query_selector('.slide')
                if slide_el:
                    img_bytes = slide_el.screenshot(type='png')
                else:
                    img_bytes = page.screenshot(type='png', full_page=False)
                slide_images.append(img_bytes)

            browser.close()
    except Exception as e:
        print(f"[PPTX Export] Playwright screenshot failed: {e}")

    # Build PPTX
    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h

    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        if i < len(slide_images):
            # Add the screenshot as a full-slide image
            from io import BytesIO
            img_stream = BytesIO(slide_images[i])
            slide.shapes.add_picture(img_stream, Emu(0), Emu(0), slide_w, slide_h)
        else:
            # Fallback: text-only slide
            title = slide_data.get('title', f'شريحة {i+1}')
            _add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1.5),
                         title, font_size=36, bold=True, font_name=font_family)

    safe_name = ''.join(c for c in project_name if c.isalnum() or c in '-_ ')[:50].strip() or 'presentation'
    output_path = os.path.join(output_dir, f"{safe_name}_{int(time.time())}.pptx")
    prs.save(output_path)
    return output_path


def _add_textbox(slide, left, top, width, height, text, font_size=18, color=None, bold=False, alignment=PP_ALIGN.RIGHT, font_name='The Sans Arabic'):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = _strip_icons(text)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if color:
        p.font.color.rgb = color
    return txBox
